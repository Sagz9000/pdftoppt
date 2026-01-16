import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
import io

def pt_to_emu(pt):
    return int(pt * 12700)

def convert_pdf_to_pptx(pdf_path: str, pptx_path: str):
    """
    Convert a PDF file to a PowerPoint presentation with layout preservation.
    """
    doc = fitz.open(pdf_path)
    prs = Presentation()
    
    # Remove default slide
    if len(prs.slides) > 0:
        # python-pptx starts with one slide? No, usually empty or one title slide.
        # Actually it starts with 0 slides if we don't add one, but let's be safe
        pass

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        rect = page.rect
        width_pt, height_pt = rect.width, rect.height
        
        # Set slide size (Note: this sets it for the whole presentation, usually based on first page)
        if page_num == 0:
            prs.slide_width = pt_to_emu(width_pt)
            prs.slide_height = pt_to_emu(height_pt)
        
        # Create blank slide
        slide_layout = prs.slide_layouts[6] # 6 is blank
        slide = prs.slides.add_slide(slide_layout)
        

        # Extract Background Drawings (Rectangles/Shapes)
        # This helps with visual layering (e.g. colored backgrounds)
        drawings = page.get_drawings()
        for draw in drawings:
            # We only handle simple filled rects for now to act as background
            # 'items' list contains drawing commands. 're' = rect, 'f' = fill
            # structure: {'items': [('re', Rect(x0,y0,x1,y1)), ...], 'fill': (r,g,b), ...}
            if draw['fill'] and draw['items']:
                rect = draw['rect']
                # Add shape
                shape = slide.shapes.add_shape(
                    1, # MSO_SHAPE.RECTANGLE
                    pt_to_emu(rect.x0), 
                    pt_to_emu(rect.y0), 
                    pt_to_emu(rect.width), 
                    pt_to_emu(rect.height)
                )
                shape.line.fill.background() # No border
                fill = shape.fill
                fill.solid()
                # draw['fill'] is (r, g, b, [a]) 0-1 floats usually? check PyMuPDF docs. 
                # PyMuPDF 'fill' is tuple of floats 0-1 if mapped? No, it's whatever the PDF space is.
                # Usually it returns (0.2, 0.4, 0.1) etc.
                try:
                    # Convert 0-1 float to 0-255 int
                    r, g, b = [int(x * 255) for x in draw['fill'][:3]]
                    fill.fore_color.rgb = RGBColor(r, g, b)
                except:
                    pass

        # Extract images
        image_list = page.get_images()
        for img_index, img in enumerate(image_list):
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                
                rects = page.get_image_rects(xref)
                if rects:
                    for img_rect in rects:
                        image_stream = io.BytesIO(image_bytes)
                        slide.shapes.add_picture(
                            image_stream, 
                            pt_to_emu(img_rect.x0), 
                            pt_to_emu(img_rect.y0), 
                            width=pt_to_emu(img_rect.width), 
                            height=pt_to_emu(img_rect.height)
                        )
            except Exception as e:
                print(f"Error adding image: {e}")

        # Extract Text
        # "dict" flag extracts text by blocks -> lines -> spans
        blocks = page.get_text("dict")["blocks"]
        for block in blocks:
            if "lines" in block:  # Text block
                bbox = block["bbox"]
                x0, y0, x1, y1 = bbox
                w = x1 - x0
                h = y1 - y0
                
                # Heuristic: Add width buffer to prevent aggressive wrapping due to font metric differences
                # PPT fonts often take more space than PDF embedded fonts
                width_buffer = 1.1 
                safe_w = w * width_buffer

                # Create text box
                txBox = slide.shapes.add_textbox(
                    pt_to_emu(x0), 
                    pt_to_emu(y0), 
                    pt_to_emu(safe_w), 
                    pt_to_emu(h) # Height matters less if not auto-fitting
                )
                tf = txBox.text_frame
                tf.word_wrap = True 
                
                first_paragraph = True
                
                # Iterate lines
                for line in block["lines"]:
                    # Create new paragraph for each line to strictly enforce vertical spacing?
                    # Or assume flow? "blocks" usually group paragraphs.
                    # To minimize overlap, let's trust the "lines" structure.
                    
                    if first_paragraph:
                        p = tf.paragraphs[0]
                        first_paragraph = False
                    else:
                        p = tf.add_paragraph()

                    # Spans within a line
                    for span in line["spans"]:
                        run = p.add_run()
                        run.text = span["text"]
                        
                        # Font size
                        run.font.size = Pt(span["size"])
                        
                        # Styles
                        if "Bold" in span["font"] or "Med" in span["font"]: # partial match
                            run.font.bold = True
                        if "Italic" in span["font"]:
                            run.font.italic = True
                            
                        # Color
                        # span["color"] is integer sRGB in "dict" mode output (e.g. 0xRRGGBB)
                        try:
                            c = span["color"]
                            # Convert integer to RGB
                            r = (c >> 16) & 0xFF
                            g = (c >> 8) & 0xFF
                            b = c & 0xFF
                            run.font.color.rgb = RGBColor(r, g, b)
                        except:
                            pass

    prs.save(pptx_path)
